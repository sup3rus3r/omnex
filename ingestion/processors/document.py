"""
Omnex — Document Processor
Extracts text from: PDF, DOCX, TXT, MD, HTML, XLSX, PPTX
Returns raw text ready for chunking.
"""

from __future__ import annotations

import io
from pathlib import Path


def extract(path: Path) -> str:
    """
    Extract text from a document file.

    Returns:
        Extracted text as a single string.
        Empty string if extraction fails or file has no text.
    """
    suffix = path.suffix.lower()

    extractors = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".doc":  _extract_docx,
        ".txt":  _extract_text,
        ".md":   _extract_text,
        ".html": _extract_html,
        ".htm":  _extract_html,
        ".xlsx": _extract_xlsx,
        ".xls":  _extract_xlsx,
        ".pptx": _extract_pptx,
        ".ppt":  _extract_pptx,
        ".json": _extract_text,
        ".xml":  _extract_html,
        ".csv":  _extract_text,
    }

    extractor = extractors.get(suffix, _extract_text)
    try:
        return extractor(path).strip()
    except Exception as e:
        # Non-fatal — log and return empty
        import logging
        logging.getLogger(__name__).warning(f"Extraction failed for {path}: {e}")
        return ""


# ── PDF ───────────────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except ImportError:
        # Fallback to PyPDF2
        import PyPDF2
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
        return "\n\n".join(pages)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n\n".join(paragraphs)


# ── Plain text / Markdown / CSV / JSON ───────────────────────────────────────

def _extract_text(path: Path) -> str:
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    # Binary fallback — extract printable chars
    raw = path.read_bytes()
    return raw.decode("utf-8", errors="replace")


# ── HTML / XML ────────────────────────────────────────────────────────────────

def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup
    html = _extract_text(path)
    soup = BeautifulSoup(html, "lxml")
    # Remove script and style
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    # Clean up excessive whitespace
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)
